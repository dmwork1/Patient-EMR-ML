global calendar, Presentation, Inches, Pt, ChartData, CategoryChartData, XL_CHART_TYPE, MSO_THEME_COLOR, Counter, namedtuple, np, pd, csv, io, os, survey, demographics, re
import calendar
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from collections import Counter
from collections import namedtuple
import numpy as np
import pandas as pd
import csv
import io
import os
import survey
import demographics
import re
# Variable Declarations

debug = 1

SLD_LAYOUT_TITLE_AND_CONTENT = 1

num_ppts = 3

survey_stats_rows = ['Survey_question_number', 'Survey_question', 'Number of responders', 'Number of non-responders', 'Mean', 'Median', 'Mode', 'Standard Deviation', 'Skewness', 'Kurtosis', 'IQR', 'Maximum', 'Minimum']

demo_rows = ['case_length', 'facility_name', 'physical_status', 'place_of_service', 'surg_npi', 'asa_desc', 'cpt', 'cpt_cat', 'emergent', 'age', 'gender', 'anes_npi', 'nurse_npi', 'usap_region']
demo_ratio_stats_rows = ['Count', 'NanNumber', 'Mean', 'Median', 'Mode', 'Std', 'Skewness', 'Kurtosis', 'IQR', 'Maximum', 'Minimum']
demo_nominal_stats_rows = ['Count', 'Mode']

# Function Declarations
def pptx_demo_ratio_table_format(slide, num_rows, num_cols):
    # ---add table to slide---#
    # x, y are positioning on slide
    # cx, cy denote table size
    shapes = slide.shapes
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(9.5), Inches(5.5)
    #Add two rows for heading and subheading minus one because question goes in slide title
    shape = slide.shapes.add_table(num_rows+1, num_cols+1, x, y, cx, cy) # first 2 args set rows, colss
    table = shape.table
    shape.has_table

    # Merge top cell
    cell = table.cell(0,0)
    other_cell = table.cell(0,(num_cols+1)-1)
    cell.merge(other_cell)

    # Change formatting
    mycols = table.columns
    mycols[0].width = Inches(5.5)
    for i in range(1,len(mycols)):
        mycols[i].width = Inches(1)
    
    myrows = table.rows
    myrows[0].height = Inches(.75)
    for i in range(1, len(myrows)):
        myrows[i].height = Inches(.25)
    return slide, shape, table, mycols, myrows

def pptx_demo_nominal_table_format(slide, num_rows, num_cols):
    # ---add table to slide---#
    # x, y are positioning on slide
    # cx, cy denote table size
    shapes = slide.shapes
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(9.5), Inches(5.5)
    #Add one rows for heading + another row for "others" row
    # Num cols + 1 accounts for num months plus 
    shape = slide.shapes.add_table(num_rows+2, num_cols+1, x, y, cx, cy) # first 2 args set rows, colss
    table = shape.table
    shape.has_table

    # Merge top cell
    cell = table.cell(0,0)
    other_cell = table.cell(0, (num_cols+1)-1)
    cell.merge(other_cell)

    # Change formatting
    mycols = table.columns
    mycols[0].width = Inches(5.5)
    for i in range(1,len(mycols)):
        mycols[i].width = Inches(1)
    
    myrows = table.rows
    myrows[0].height = Inches(.75)
    for i in range(1, len(myrows)):
        myrows[i].height = Inches(.25)
    return slide, shape, table, mycols, myrows
    
def pptx_survey_table_format(slide, num_rows, num_cols):
    # ---add table to slide---#
    # x, y are positioning on slide
    # cx, cy denote table size
    shapes = slide.shapes
    x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(9.5), Inches(1)#Inches(5.5)
    #Add two rows for heading and subheading minus one because question goes in slide title
    shape = slide.shapes.add_table(num_rows, num_cols+1, x, y, cx, cy) # first 2 args set rows, colss
    table = shape.table
    shape.has_table

    # Merge top cell
    cell = table.cell(0,0)
    other_cell = table.cell(0,(num_cols+1)-1)
    cell.merge(other_cell)

    # Change formatting
    mycols = table.columns
    mycols[0].width = Inches(4)
    for i in range(1,len(mycols)):
        mycols[i].width = Inches(1)
    
    myrows = table.rows
    #myrows[0].height = Inches(.75)
    for i in range(1, len(myrows)):
        myrows[i].height = Inches(.15)
    return slide, shape, table, mycols, myrows

def open_pptx(year, i):
    # Open 
    prs = Presentation('USAP Statistics ' + str(year+i) + ' orig.pptx')    
    slide_layout = prs.slide_layouts[5]
    return prs, slide_layout
    
def write_pptx_demo_ratio_data(prs, slide_layout, init_index, demo_stats, survey_file):    
    #
    print('Placeholder')
    
    return prs, idx
    
def write_pptx_demo_nominal_data(prs, slide_layout, init_index, demo_stats, survey_file):    
    #
    print('Placeholder')
    
    return prs, idx
    
def write_pptx_survey_data(prs, slide_layout, init_index, num_months, survey_stats, survey_file):

    # Will need to pass slide_layout to stats function
    # so make sure these function files contain the imports (as global vars)
    #print('gotcha')
    #print('np.shape(survey_stats)[1] =',str(np.shape(survey_stats)[1]))
    for y in range(0, 22):#np.shape(survey_stats)[1]):    # Number of Questions 
        print('fix this above!!! hardcoding')
        # Add slide
        slide = prs.slides.add_slide(slide_layout)
        # Select appropriate slide
        try:
            idx = init_index+(np.shape(survey_stats)[1]*0+y)+1 #indexing for appropriate slide, +1 accounts for title slide
        except:
            idx = init_index+y+1
            slide.shapes.title.text = 'Garbage Slide'
            slide = prs.slides[idx]
            
            continue
        slide = prs.slides[idx]
        #Modify slide title
        slide.shapes.title.text = 'Question ' + str(survey_stats[0][y][0])
        print('New Title is:')
        print(slide.shapes.title.text)
        #Format survey slide
        # Currently 4 months at a time
        slide, shape, table, mycols, myrows = pptx_survey_table_format(slide, len(survey_stats_rows), np.shape(survey_stats)[0])
        #print top 10 most common - since this list can change, only one month at a time
        #slide, shape, table, mycols, myrows = pptx_demo_nominal_table_format(slide, 10, 1, np.shape(survey_stats)[0])
        # Don't include "Count" for ratio data
        #slide, shape, table, mycols, myrows = pptx_demo_ratio_table_format(slide, len(demo_ratio_stats_rows)-1, np.shape(survey_stats)[0])
        # Access cells and write
        tabletitlecell = table.cell(0,0)
        tabletitlecell.text = survey_stats[0][y][1] #hardcoded
        paragraph = tabletitlecell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(16)
        print('My rows length =', len(myrows))
        print('My cols length =', len(mycols))
        print('Range(2, len(myrows)) = ',str(range(2,len(myrows))))
        print('Length of survey_stats_rows = ',len(survey_stats_rows))
        print('Survey stats rows: ',survey_stats_rows)
        cell = table.cell(1,0)
        cell.text = 'Statistics'
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)        
        for x in range(0, np.shape(survey_stats)[0]): 
            month_index = np.int32((re.findall(r'(\d{4})\_(\d{2})\_\w+\.\w{3}', survey_file[x]))[0][1])
            table.cell(1,x+1).text = calendar.month_abbr[month_index] 
            paragraph = table.cell(1,x+1).text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            smile = 0
            for k in range(0, len(myrows)-2):
                print('Len(myrows)-2 = ',len(myrows)-2)
                #table.cell(k+2,0).text = survey_stats_rows[k+2]
                if x == 0:
                    p = table.cell(k+2,0).text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = survey_stats_rows[k+2]
                    font = run.font
                    font.size = Pt(14)
                    font.bold = False
                font.italic = None  # cause value to be inherited from theme
                p1 = table.cell(k+2,x+1).text_frame.paragraphs[0]               
                run1 = p1.add_run()                
                print('Np.shape(survey_stats) ==',str(np.shape(survey_stats)))
                print('Type survey_stats[x][y][k+2]',str(type(survey_stats[x][y][2])))
                if smile >= 2:
                    run1.text = "{:.3f}".format(survey_stats[x][y][k+2])
                else:
                    run1.text = str(survey_stats[x][y][k+2])                
                font1 = run1.font                
                font1.size = Pt(14)                
                smile += 1
            print('checkpoint 3')
            # Debug code
            if debug:
                shape.has_table
                table = shape.table
                table
    return prs, idx
def close_pptx(year, i, prs):
    # Finished    
    prs.save('USAP Statistics ' + str(year+i) + ' automated.pptx')
# Open Presentation - 1 slide PPT already made
#for j in range(num_ppts): # there should be three, this loop exists already as I loop in USAP statistics
   

    